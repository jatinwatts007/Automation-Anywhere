from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
def add_text(input2):
    prs=Presentation("C:/Users/mayan/OneDrive/Desktop/Juypter/INT346/test_4.pptx")
    lyt=prs.slide_layouts[0] # choosing a slide layout
    slide=prs.slides.add_slide(lyt) # adding a slide
    title=slide.shapes.title # assigning a title
    subtitle=slide.placeholders[1] # placeholder for subtitle
    title.text=input2[0] # title
    subtitle.text=input2[1] # subtitle
    prs.save("C:/Users/mayan/OneDrive/Desktop/Juypter/INT346/test_4.pptx") # saving file

def add_graph(input1):
    # create presentation with 1 slide ------
    prs = Presentation("C:/Users/mayan/AppData/Local/Temp/example.pptx")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # define chart data ---------------------
    chart_data = CategoryChartData()
    chart_data.categories = ['East', 'West', 'Midwest']
    chart_data.add_series('Series 1', (input1[0], input1[1], input1[2]))

    # add chart to slide --------------------
    title=slide.shapes.title
    title.text="Bar Graph" # title
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    prs.save("C:/Users/mayan/AppData/Local/Temp/example.pptx")

    
def add_image(input3):
    prs = Presentation("C:/Users/mayan/AppData/Local/Temp/example.pptx")
    layout8 = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout8)
    title = slide.shapes.title.text = input3[0]
    sub = slide.placeholders[2].text = input3[1]
    placeholder = slide.placeholders[1]
 
    # Calculate the image size of the image
    im = Image.open(input3[2])
    width, height = im.size
 
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width
 
    # Insert the picture
    placeholder = placeholder.insert_picture(input3[2])
 
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
 
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side
    prs.save("C:/Users/mayan/AppData/Local/Temp/example.pptx")
    
    
def add_shape():
    prs = Presentation("C:/Users/mayan/AppData/Local/Temp/example.pptx")
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Adding an AutoShape'

    left = Inches(0.93)  # 0.93" centers this overall set of shapes
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)

    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = 'Step 1'

    left = left + width - Inches(0.4)
    width = Inches(2.0)  # chevrons need more width for visual balance

    for n in range(2, 6):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = 'Step %d' % n
        left = left + width - Inches(0.4)

    prs.save("C:/Users/mayan/AppData/Local/Temp/example.pptx")
